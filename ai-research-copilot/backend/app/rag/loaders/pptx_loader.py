"""
PowerPoint document loader.

Extracts text from PowerPoint (.pptx) files using python-pptx.
"""

import logging
from io import BytesIO

from pptx import Presentation

from app.rag.loaders.base import BaseDocumentLoader

logger = logging.getLogger(__name__)


class PPTXLoader(BaseDocumentLoader):
    """Load and extract text from PowerPoint presentations."""

    async def load(self, content: bytes, filename: str) -> str:
        """Extract text from a PPTX file.

        Iterates through all slides and shapes, extracting text from
        text boxes, placeholders, and tables.

        Args:
            content: Raw PPTX bytes.
            filename: Original filename for logging.

        Returns:
            Extracted text with slides separated by double newlines.

        Raises:
            ValueError: If the PPTX cannot be parsed.
        """
        try:
            prs = Presentation(BytesIO(content))
        except Exception as exc:
            raise ValueError(f"Failed to parse PPTX '{filename}': {exc}") from exc

        slide_texts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            shape_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            shape_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            shape_texts.append(row_text)
            if shape_texts:
                slide_texts.append("\n".join(shape_texts))

        full_text = "\n\n".join(slide_texts)

        if not full_text.strip():
            raise ValueError(f"PPTX '{filename}' contains no extractable text")

        logger.info("Extracted text from %d slides in '%s'", len(slide_texts), filename)
        return full_text
